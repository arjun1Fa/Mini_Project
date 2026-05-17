from pptx import Presentation
from pptx.util import Inches, Pt

def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    
    tf = content_placeholder.text_frame
    tf.text = content
    # Make font slightly smaller so it fits well
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)

prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0] # Title layout
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "NutriLens (or NutriVision)"
slide.placeholders[1].text = "Precision AI-Powered Nutrition Tracking for Indian Cuisine\n\nPresenter: [Your Name/Team Name]"

# Slide 2
add_slide(prs, "The Flaw in Modern Calorie Counting", 
"• Blind Guessing: Most apps force users to manually guess the weight of their food, leading to massive errors.\n"
"• Western Bias: Existing AI scanners fail at complex ethnic cuisines (e.g., Aviyal, Puttu, Parotta).\n"
"• Friction: Manually logging every ingredient is tedious and causes users to quit.")

# Slide 3
add_slide(prs, "NutriLens - See Your Nutrition", 
"• Volumetric AI: Calculates exactly how much food is on the plate using advanced depth estimation.\n"
"• Tailored for Indian Food: Custom 20-class taxonomy for authentic Indian and Kerala cuisines.\n"
"• Seamless UX: Snap a photo -> Get exact calories, protein, carbs, and fat instantly.")

# Slide 4
add_slide(prs, "A Robust, Modern Tech Stack", 
"• Frontend: Flutter & Riverpod for a beautiful, cross-platform fluid UX.\n"
"• Backend: Highly scalable Python/Flask REST API handling ML inference.\n"
"• Database: Supabase (PostgreSQL) for secure auth, RLS, and history storage.\n"
"• Cloud Ready: Dockerized backend designed for GCP Cloud Run.")

# Slide 5
add_slide(prs, "Dual-Model Computer Vision", 
"• Model 1: Instance Segmentation (YOLOv8-seg)\n"
"  - Identifies food type (e.g., 'Idli', 'Chicken Curry').\n"
"  - Draws precise pixel-perfect boundaries around every piece of food.\n\n"
"• Model 2: Monocular Depth (Depth Anything V2)\n"
"  - Calculates the 3D topology of the plate.\n"
"  - Determines exactly how 'tall' or 'thick' the food is in centimeters.")

# Slide 6
add_slide(prs, "How We Calculate Weight", 
"1. Area Extraction: YOLOv8 gives the 2D surface area of the food.\n"
"2. Height Extraction: Depth Anything V2 gives the 3D height.\n"
"3. Volume Calculation: Area × Height = Exact Volume (cm³).\n"
"4. Density Application: Cross-reference volume with custom shape_density_table.\n"
"5. Result: Exact gram weight (Volume × Density = Weight).")

# Slide 7
add_slide(prs, "Powered by Standardized Data", 
"• IFCT 2017 Integration: Database strictly mapped to the Indian Food Composition Tables (IFCT 2017).\n"
"• Smart Fallbacks: Manual-entry search instantly queries the database to auto-fill IFCT macro values.\n"
"• Macro Tracking: Real-time tracking of Protein, Carbs, Fats, and Fiber against personalized goals.")

# Slide 8
add_slide(prs, "Premium Aesthetics & Gamification", 
"• Dynamic UI: Glassmorphism, smooth micro-animations, and curated palettes.\n"
"• Analytics Dashboard: Visualizing weekly calorie trends with interactive charts.\n"
"• Offline Resilience: Gracefully handles network issues, queuing manual entries seamlessly.")

# Slide 9
add_slide(prs, "Seeing is Believing (Demo)", 
"• Showcase 1: Taking a picture of a complex meal (e.g., Parotta and Curry) and showing the AI breakdown.\n"
"• Showcase 2: The history screen, showing rich emoji mappings and daily goal progress.\n\n"
"(Live Demo or Video Integration here)")

# Slide 10
add_slide(prs, "What's Next for NutriLens?", 
"• Expanded Taxonomy: Scaling from 20 classes to 100+ regional Indian dishes.\n"
"• Edge AI: Moving lighter versions of the models directly to the phone (TFLite) for zero-latency, offline inference.\n"
"• Wearable Integration: Syncing caloric expenditure with Apple Health and Google Fit.")

# Slide 11
add_slide(prs, "Thank You", 
"Q&A\n\nOpen the floor to the judges for technical or design-related questions.")

prs.save('NutriLens_Presentation.pptx')
print("Successfully generated NutriLens_Presentation.pptx")
